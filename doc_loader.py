"""
doc_loader.py — 文档加载、清洗与切分

支持格式: PDF / TXT / Markdown / PPTX
清洗管线: 页眉页脚移除 → 断行合并 → 重复去重 → 特殊字符清洗 → 空白归一化
切分策略:
  - recursive: 递归字符切分（默认，通用）
  - character: 固定字符切分
  - markdown: 按标题切分（仅 md 文件）

提供切分效果对比工具，用于可视化不同策略的差异。
"""
import os
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Literal

try:
    from langchain_core.documents import Document
    from langchain_text_splitters import (
        RecursiveCharacterTextSplitter,
        CharacterTextSplitter,
        MarkdownHeaderTextSplitter,
    )
    HAS_LANGCHAIN_SPLITTERS = True
except ImportError:
    HAS_LANGCHAIN_SPLITTERS = False

from config import get_config


@dataclass
class LoadResult:
    """文档加载结果"""
    documents: list = field(default_factory=list)
    total_files: int = 0
    total_chunks: int = 0
    failed_files: list = field(default_factory=list)


class TextCleaner:
    """
    文档数据清洗器

    清洗管线（按顺序执行）:
      1. remove_headers_footers  — 移除页眉页脚、页码
      2. merge_broken_lines    — 合并 PDF 断行破碎的句子
      3. deduplicate_text      — 去除重复段落（页眉重复等）
      4. clean_special_chars   — 清理特殊字符、乱码
      5. normalize_whitespace  — 空白归一化

    使用:
      cleaner = TextCleaner(config)
      cleaned_text = cleaner.clean(raw_text)
    """

    # 页码/页眉页脚常见模式
    HEADER_FOOTER_PATTERNS = [
        re.compile(r'^[-—]\s*第\s*\d+\s*页\s*[-—]?\s*$'),
        re.compile(r'^第\s*\d+\s*页\s*/\s*共\s*\d+\s*页\s*$'),
        re.compile(r'^\d+\s*/\s*\d+\s*$'),
        re.compile(r'^Page\s+\d+\s*(of\s+\d+)?\s*$', re.IGNORECASE),
        re.compile(r'^[-=]{3,}\s*$'),
        re.compile(r'^\d+\s*$'),
        re.compile(r'第\s*\d+\s*页\s*/\s*共\s*\d+\s*页'),
        re.compile(r'^.+?\s+第\s*\d+\s*页\s*/\s*共\s*\d+\s*页\s*$'),
        re.compile(r'^.+?\s+\d+\s*/\s*\d+\s*$'),
    ]

    # 特殊字符/乱码模式
    SPECIAL_CHAR_PATTERNS = [
        (re.compile(r'\u3000'), ' '),
        (re.compile(r'&nbsp;'), ' '),
        (re.compile(r'&[a-zA-Z]+;'), ' '),
        (re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f]'), ''),
        (re.compile(r'[\ufffe\uffff]'), ''),
    ]

    def __init__(self, config=None):
        if config is not None and hasattr(config, 'cleaner'):
            self.cfg = config.cleaner
        elif config is not None and hasattr(config, 'enabled'):
            self.cfg = config
        else:
            from config import CleanerConfig
            self.cfg = CleanerConfig()

    def clean(self, text: str, source: str = "") -> str:
        """执行完整清洗管线"""
        if not text or not text.strip():
            return ""
        if not self.cfg.enabled:
            return self._normalize(text)

        # 1. 特殊字符清洗（先做，避免干扰后续步骤）
        if self.cfg.clean_special_chars:
            text = self._clean_special_chars(text)

        # 2. 页眉页脚移除
        if self.cfg.remove_headers_footers:
            text = self._remove_headers_footers(text)

        # 3. 断行合并
        if self.cfg.merge_broken_lines:
            text = self._merge_broken_lines(text)

        # 4. 重复去重
        if self.cfg.deduplicate_text:
            text = self._deduplicate(text)

        # 5. 空白归一化
        if self.cfg.normalize_whitespace:
            text = self._normalize(text)

        return text

    def _remove_headers_footers(self, text: str) -> str:
        """移除页眉页脚"""
        lines = text.split('\n')
        kept = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                kept.append(line)
                continue
            is_noise = False
            for pattern in self.HEADER_FOOTER_PATTERNS:
                if pattern.match(stripped):
                    is_noise = True
                    break
            if not is_noise:
                kept.append(line)
        return '\n'.join(kept)

    def _merge_broken_lines(self, text: str) -> str:
        """
        合并 PDF 断行破碎的句子。

        PDF 提取时，一句话经常被断成多行（如自动换行导致）。
        规则: 不以句末标点结尾的行，与下一行合并。
        """
        lines = text.split('\n')
        merged = []
        sentence_endings = set('。！？.!?；;：:')

        i = 0
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            if not stripped:
                merged.append(line)
                i += 1
                continue

            # 如果这一行不以句末标点结尾，且不是空行，则尝试合并
            while (i + 1 < len(lines)
                   and stripped
                   and stripped[-1] not in sentence_endings
                   and len(stripped) > self.cfg.min_line_length):
                next_line = lines[i + 1].strip()
                if not next_line:
                    break
                stripped = stripped + next_line
                i += 1

            merged.append(stripped)
            i += 1

        return '\n'.join(merged)

    def _deduplicate(self, text: str) -> str:
        """
        去除重复段落。

        检测被重复的文本块（通常是每页的页眉/页脚），
        使用字符集合 Jaccard 相似度进行判定。
        """
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if len(paragraphs) <= 1:
            return text

        unique = []
        for para in paragraphs:
            is_dup = False
            for existing in unique:
                if self._jaccard_similarity(para, existing) >= self.cfg.similarity_threshold:
                    is_dup = True
                    break
            if not is_dup:
                unique.append(para)

        return '\n\n'.join(unique)

    def _jaccard_similarity(self, text1: str, text2: str) -> float:
        """计算两段文本的 Jaccard 相似度"""
        set1 = set(text1)
        set2 = set(text2)
        if not set1 or not set2:
            return 0.0
        intersection = set1 & set2
        union = set1 | set2
        return len(intersection) / len(union)

    def _clean_special_chars(self, text: str) -> str:
        """清理特殊字符和乱码"""
        for pattern, replacement in self.SPECIAL_CHAR_PATTERNS:
            text = pattern.sub(replacement, text)
        return text

    def _normalize(self, text: str) -> str:
        """空白归一化"""
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = '\n'.join(line.rstrip() for line in text.split('\n'))
        return text.strip()


class DocumentLoader:
    """多格式文档加载器"""

    SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".pptx"}

    def __init__(self, config=None):
        self.cfg = config or get_config()
        self.splitter_cfg = self.cfg.splitter
        self.cleaner = TextCleaner(self.cfg)

    def load_directory(self, dir_path: str | Path) -> LoadResult:
        """加载目录下所有支持的文档"""
        dir_path = Path(dir_path)
        if not dir_path.exists():
            return LoadResult()

        result = LoadResult()

        for root, _, files in os.walk(dir_path):
            for fname in files:
                fpath = Path(root) / fname
                if fpath.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                    result.total_files += 1
                    try:
                        docs = self.load_file(fpath)
                        result.documents.extend(docs)
                    except Exception as e:
                        result.failed_files.append((str(fpath), str(e)))

        result.total_chunks = len(result.documents)
        return result

    def load_file(self, file_path: str | Path) -> list:
        """加载单个文件，清洗后切分"""
        file_path = Path(file_path)
        ext = file_path.suffix.lower()

        # 1. 读取原始文本
        if ext in (".txt", ".md"):
            raw_text = file_path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".pdf":
            raw_text = self._load_pdf(file_path)
        elif ext == ".pptx":
            raw_text = self._load_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")

        # 2. 数据清洗
        cleaned_text = self.cleaner.clean(raw_text, source=str(file_path))

        # 3. 切分
        docs = self._split_text(cleaned_text, str(file_path), ext)
        return docs

    def _split_text(self, text: str, source: str, ext: str) -> list:
        """根据策略切分文本"""
        if not HAS_LANGCHAIN_SPLITTERS:
            # Fallback：简单按段落切分
            return self._simple_split(text, source)

        strategy = self.splitter_cfg.strategy

        if strategy == "markdown" and ext == ".md":
            splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=[
                    ("#", "Header 1"),
                    ("##", "Header 2"),
                    ("###", "Header 3"),
                ]
            )
            docs = splitter.split_text(text)
            for d in docs:
                d.metadata["source"] = source
            return docs

        elif strategy == "character":
            splitter = CharacterTextSplitter(
                chunk_size=self.splitter_cfg.chunk_size,
                chunk_overlap=self.splitter_cfg.chunk_overlap,
                separator="\n",
            )
        else:  # recursive
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.splitter_cfg.chunk_size,
                chunk_overlap=self.splitter_cfg.chunk_overlap,
                separators=["\n\n", "\n", "。", "！", "？", ";", ".", " "],
            )

        docs = splitter.create_documents([text], metadatas=[{"source": source}])
        return docs

    def _simple_split(self, text: str, source: str) -> list:
        """简单切分（无依赖时的 fallback）"""
        if not HAS_LANGCHAIN_SPLITTERS:
            # 纯 Python 简单切分
            chunks = []
            chunk_size = self.splitter_cfg.chunk_size
            overlap = self.splitter_cfg.chunk_overlap
            paragraphs = text.split("\n\n")

            current = ""
            for para in paragraphs:
                if len(current) + len(para) > chunk_size and current:
                    chunks.append(self._make_doc(current, source, len(chunks)))
                    # 保留 overlap 长度的 overlap
                    current = current[-overlap:] + para + "\n\n"
                else:
                    current += para + "\n\n"

            if current.strip():
                chunks.append(self._make_doc(current, source, len(chunks)))

            return chunks
        return []

    def _make_doc(self, content: str, source: str, idx: int):
        """构造 Document 对象"""
        if HAS_LANGCHAIN_SPLITTERS:
            from langchain_core.documents import Document
            return Document(page_content=content.strip(), metadata={"source": source, "page": idx + 1})
        else:
            return {"page_content": content.strip(), "metadata": {"source": source, "page": idx + 1}}

    def _load_pdf(self, file_path: Path) -> str:
        """
        加载 PDF 文件（智能路由）

        自动判断每页类型并路由:
          - 原生 PDF  → 直接文本抽取
          - 扫描件    → OCR 识别
          - 图表/表格 → 多模态视觉理解
        """
        try:
            from pdf_processor import PDFProcessor
            processor = PDFProcessor(self.cfg)
            return processor.process_to_text(file_path)
        except ImportError:
            # 降级到 pypdf
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(file_path))
                texts = []
                for i, page in enumerate(reader.pages):
                    t = page.extract_text() or ""
                    texts.append(f"--- 第 {i+1} 页 ---\n{t}")
                return "\n\n".join(texts)
            except ImportError:
                raise ImportError(
                    "加载 PDF 需要 PyMuPDF 或 pypdf 库。\n"
                    "请运行: pip install PyMuPDF"
                )

    def _load_pptx(self, file_path: Path) -> str:
        """加载 PPTX 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    texts.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(slide_text))
            return "\n\n".join(texts)
        except ImportError:
            raise ImportError(
                "加载 PPTX 需要 python-pptx 库。\n请运行: pip install python-pptx"
            )

    def compare_split_strategies(self, file_path: str | Path) -> dict:
        """
        对比三种切分策略的效果。

        Returns:
            dict: {strategy: {"chunks": int, "avg_size": float, "sample": str}}
        """
        file_path = Path(file_path)
        text = ""
        ext = file_path.suffix.lower()

        if ext in (".txt", ".md"):
            text = file_path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".pdf":
            text = self._load_pdf(file_path)

        results = {}

        if not HAS_LANGCHAIN_SPLITTERS:
            return {"recursive": {"chunks": 0, "avg_size": 0, "sample": "未安装 langchain"}}

        chunk_size = self.splitter_cfg.chunk_size
        chunk_overlap = self.splitter_cfg.chunk_overlap

        for strategy in ["recursive", "character", "markdown"]:
            if strategy == "markdown" and ext != ".md":
                continue

            if strategy == "recursive":
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                    separators=["\n\n", "\n", "。", " "],
                )
            elif strategy == "character":
                splitter = CharacterTextSplitter(
                    chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                    separator="\n",
                )
            else:
                splitter = MarkdownHeaderTextSplitter(
                    headers_to_split_on=[("#", "H1"), ("##", "H2"), ("###", "H3")]
                )

            if strategy == "markdown":
                docs = splitter.split_text(text)
            else:
                docs = splitter.create_documents([text])

            sizes = [len(d.page_content) for d in docs]
            results[strategy] = {
                "chunks": len(docs),
                "avg_size": round(sum(sizes) / len(sizes), 1) if sizes else 0,
                "max_size": max(sizes) if sizes else 0,
                "min_size": min(sizes) if sizes else 0,
                "sample": docs[0].page_content[:100] if docs else "",
            }

        return results
